# -*- coding: utf-8 -*-
"""
PowerPoint文件适配器

支持 .ppt 和 .pptx 文件的解析。
"""

from pathlib import Path

from .base_adapter import BaseFileAdapter


class PptFileAdapter(BaseFileAdapter):
    """PowerPoint文件适配器"""

    def __init__(self, file_path: str, max_file_size: int = 30 * 1024 * 1024):
        super().__init__(file_path, max_file_size)
        self._file_content: str | None = None
        self._slides_content: list[str] | None = None

    def _get_file_description(self) -> str:
        """获取文件描述"""
        return "PowerPoint Presentation"

    async def _load_presentation(self) -> list[str] | None:
        """加载演示文稿内容"""
        if self._slides_content is not None:
            return self._slides_content

        from aiofiles import open as aio_open
        import aiofiles

        stat = await aiofiles.os.stat(self.file_path)
        if stat.st_size > self.max_file_size:
            return None

        try:
            async with aio_open(self.file_path, "rb") as f:
                buffer = await f.read()

            # 使用 python-pptx 解析 PPTX
            from pptx import Presentation
            from io import BytesIO

            prs = Presentation(BytesIO(buffer))

            slides_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []

                # 提取文本框内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

                # 合并幻灯片文本
                full_text = "\n".join(slide_texts)
                slides_text.append(f"Slide {slide_num}:\n{full_text}")

            self._slides_content = slides_text
            return self._slides_content

        except Exception as e:
            from app.core.logger import log
            log.error(f"Error parsing PowerPoint file: {e}")
            return None

    async def get_content(self) -> str | None:
        """获取文件原始内容"""
        if self._file_content is None:
            slides = await self._load_presentation()
            if slides:
                self._file_content = "\n\n".join(slides)

        return self._file_content

    async def get_llm_content(self) -> str | None:
        """获取适合LLM处理的内容格式"""
        slides = await self._load_presentation()
        if not slides:
            return None

        file_description = f"""# PowerPoint Presentation Description

## Basic PowerPoint Presentation Information
* **Total Slides:** {len(slides)}

## Slides Information
"""

        # 为每张幻灯片生成内容
        slides_content = []
        for i, slide_text in enumerate(slides, 1):
            # 提取幻灯片内容（去掉 "Slide X:" 前缀）
            content = slide_text.split("\n", 1)[1] if "\n" in slide_text else ""
            slides_content.append(f"""### Slide {i}
```
{content}
```
""")

        return file_description + "\n".join(slides_content)

    async def get_thumbnail(self) -> str | None:
        """获取缩略图"""
        return None
